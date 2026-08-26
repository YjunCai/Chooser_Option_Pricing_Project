"""
Week 8 Final Presentation Builder
==================================
Generates Week_8_Final_Presentation.pptx (final demonstration deck) with
python-pptx. Slides cover the full 8-week arc, key results, the finalized
pricing tool, and investment insights.

Run:
    python make_presentation.py
"""

import sys
import warnings
from pathlib import Path

warnings.filterwarnings('ignore')
sys.path.insert(0, str(Path(__file__).resolve().parent))

import w8config as cfg

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ────────────────────────────────────────────────────────────────────
DARK = RGBColor(0x1F, 0x3B, 0x57)
BLUE = RGBColor(0x21, 0x66, 0xAC)
RED = RGBColor(0xD7, 0x30, 0x27)
AMBER = RGBColor(0xFD, 0xAE, 0x61)
GREY = RGBColor(0x4A, 0x4A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)

ASSETS = cfg.ASSETS_DIR


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _title_box(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = DARK
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(13); r2.font.color.rgb = GREY
    return box


def _bullets(slide, items, top=1.4, left=0.7, width=12.0, size=17, height=6.2):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run()
        r.text = ('• ' if lvl == 0 else '– ') + txt
        r.font.size = Pt(size - 2 * lvl)
        r.font.color.rgb = GREY if lvl else DARK
        p.space_after = Pt(6)
    return box


def _metric_table(slide, rows, top=1.6, col_w=None):
    nrows, ncols = len(rows), len(rows[0])
    tbl_shape = slide.shapes.add_table(nrows, ncols, Inches(0.7), Inches(top),
                                       Inches(12.1), Inches(0.5 * nrows))
    tbl = tbl_shape.table
    if col_w:
        for i, w in enumerate(col_w):
            tbl.columns[i].width = Inches(w)
    for j, cell in enumerate(rows[0]):
        c = tbl.cell(0, j)
        c.text = cell
        c.text_frame.paragraphs[0].runs[0].font.bold = True
        c.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        c.fill.solid(); c.fill.fore_color.rgb = DARK
        c.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    for i, row in enumerate(rows[1:], start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            c.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
    return tbl_shape


def _add_picture(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return None
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                    width=Inches(width) if width else None,
                                    height=Inches(height) if height else None)


def _footer(slide, idx, total):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.3), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f'{idx} / {total}'
    r.font.size = Pt(10); r.font.color.rgb = GREY


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = []
    total = 14

    # ── 1. title ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, DARK)
    box = s.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11), Inches(3.4))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '基于真实数据与机器学习的\n高级选择权定价模型'
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = 'Quantitative Research & Trading — Chooser Option Pricing (8 Weeks)'
    r2.font.size = Pt(18); r2.font.color.rgb = AMBER
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = 'Week 8 最终交付 · 可部署定价工具 + 最终报告 + 演示'
    r3.font.size = Pt(16); r3.font.color.rgb = WHITE
    slides.append(s)

    # ── 2. project overview ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '项目概述与目标', 'Project Overview & Objectives')
    _bullets(s, [
        ('选择权期权（Chooser Option）：选择日决定行权为欧式看涨/看跌，统一 K 与 T2', 0),
        ('业务价值：结构化产品、风险管理（财报/FOMC）、专有交易策略', 0),
        ('核心问题：BSM 假设恒定波动率，需以真实数据 + ML 提升定价精度', 0),
        ('5 大目标：BSM 复现验证 → 多维数据 → ML 优化波动率 → 量化改进 → 可部署工具', 0),
        ('8 周推进：数据(1-2) → BSM 基线(3-4) → 双轨 ML(5-6) → 敏感性与工具(7-8)', 0),
    ])
    _footer(s, 2, total); slides.append(s)

    # ── 3. data & features ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '数据与特征工程（第 1–2 周）', 'Data Collection & Feature Engineering')
    _bullets(s, [
        ('数据源：Yahoo Finance / Alpha Vantage / FRED；JPM 日线 + VIX + 3M 国债', 0),
        ('区间 2018-01–2024-12，清洗后 1755 行 × 20 列（插值缺失 / IQR 剔除异常）', 0),
        ('基础特征：5/21/63 日滚动波动率、日收益、股息增长、区间振幅', 0),
        ('新特征：sentiment_score(0-1)、vix_ratio、VIX-JPM 相关/交叉、利率动量', 0),
        ('自动化流水线 + GitHub Actions 调度，全部特征后视构造防前瞻', 0),
    ], top=1.5)
    _footer(s, 3, total); slides.append(s)

    # ── 4. BSM replication ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, 'BSM 模型复现与基线评估（第 3–4 周）', 'BSM Replication & Baseline')
    _bullets(s, [
        ('Rubinstein (1991) 简单选择权：P = C(S,K,T2) + P(S,Ke^{-q(T2-t1)}, t1)', 0),
        ('参数对齐论文：K = $150、T2 = 1 年、t1 = 0.5 年', 0),
        ('BSM(vol_21d) 基线（2026-07-27 实盘快照，595 合约）：', 0),
    ], top=1.5, height=2.4)
    _metric_table(s, [
        ['指标', 'Week 4 基线', '说明'],
        ['实盘 MAE ($)', '$1.44', 'BSM(vol_21d) 静态波动率'],
        ['ATM IV 溢价 (%)', '4.9', '模型 IV vs 市场 ATM IV'],
        ['OTM 看跌定价偏差 (%)', '$-65.7$', '系统低估 OTM 看跌'],
        ['合成测试集 Chooser MAE ($)', '$1.52', '2024 共 208 日'],
    ], top=3.9)
    _footer(s, 4, total); slides.append(s)

    # ── 5. ML dual-track ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '机器学习双轨架构（第 5 周）', 'Dual-Track ML Architecture')
    _bullets(s, [
        ('方法一（首选）：ML 波动率预测 + BSM 定价引擎', 0),
        ('  LSTM / RandomForest / GBDT / XGBoost / VIX-proxy（以 VIX/100 为 σ）', 1),
        ('方法二：端到端监督定价（Linear / GBDT / MLP，合约网格特征）', 0),
        ('反前瞻保障：70/15/15 时序分割、特征后视、目标错位、Scaler 训练集拟合', 0),
        ('Week 6 调优前诊断：已实现波动率目标噪声大 → 系统性加大正则', 0),
    ], top=1.5)
    _footer(s, 5, total); slides.append(s)

    # ── 6. week6 results (vol + chooser) ────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '模型训练、调优与比较（第 6 周）', 'Training, Tuning & Comparison')
    _bullets(s, [
        ('超参优化：purged 时序 CV（扩张窗口 + 21 日 embargo）+ 随机搜索 → 强正则', 0),
        ('波动率预测首次越过 persistence 门槛：GBDT-anchored 7.16% < 7.32%', 0),
        ('Chooser 定价（合成测试集）：VIX-proxy $1.17 vs BSM $1.52，改善 23%', 0),
        ('方法二调优后仍逊于静态 BSM（$3.69 / $4.61 vs $1.77）→ 方法一为最终首选', 0),
        ('SHAP：以市场隐含波动率为目标时，位置/状态特征（sentiment、vix_ratio）主导', 0),
    ], top=1.45)
    _add_picture(s, ASSETS / 'fig_w8_vol_metrics.png', 7.2, 1.9, width=5.6)
    _footer(s, 6, total); slides.append(s)

    # ── 7. live targets ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '实盘指标：三项目标全部达成', 'Live-Snapshot Targets Achieved (XGBoost)')
    _metric_table(s, [
        ['指标', 'Week 4 基线', 'Week 5 (RF)', 'Week 6 最优 (XGB)', '目标', '达成'],
        ['实盘 MAE ($)', '1.44', '1.03', '0.79', '< 1.00', '✓'],
        ['ATM IV 溢价 (%)', '4.9', '3.0', '0.76', '< 2.0', '✓'],
        ['OTM 看跌偏差 (%)', '$-65.7$', '$-51.4$', '$-29.5$', '> -30', '✓'],
        ['Chooser MAE ($)', '1.52', '1.46', '1.17', '持续下降', '✓'],
    ], top=1.8, col_w=[2.6, 2.1, 2.1, 2.3, 1.9, 1.1])
    _bullets(s, [
        ('调优后 XGBoost 预测 σ ≈ 24.4%，贴近市场 ATM IV 25.2% → IV premium 被捕获', 0),
        ('实盘与合成口径排序相反：实盘偏好 XGB（贴近市场 IV），合成偏好 VIX-proxy', 0),
    ], top=4.6)
    _footer(s, 7, total); slides.append(s)

    # ── 8. sensitivity ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '高级敏感性分析（第 7 周）', 'Extended Sensitivity Analysis')
    _bullets(s, [
        ('SHAP×Vega 分解：情感/VIX 新特征是价格第一驱动', 0),
        ('  sentiment_score 价格影响 $0.314（0.54% 均值价格），是次高特征 3 倍', 1),
        ('ATM 基座极端场景：波动率 +50% → BSM +49% / XGB +70% / VIX-proxy +81%', 0),
        ('VIX 冲击 +50%：BSM 无反应，XGB +14% / VIX-proxy +21%（特征重工程）', 0),
        ('实盘基座（S=356，深度实值）：波动率敏感性趋零（<0.3%）', 0),
    ], top=1.45)
    _add_picture(s, ASSETS / 'fig_w7_scenarios_atm.png', 7.0, 2.0, width=6.0)
    _footer(s, 8, total); slides.append(s)

    # ── 9. final tool: dashboard ────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '最终定价工具（第 8 周）：可视化仪表板', 'Final Tool — Complete Dashboard')
    _bullets(s, [
        ('双轨定价：BSM(vol_21d) vs 最优 ML（XGB / VIX-proxy / GBDT-anchored）', 0),
        ('误差条：每项定价 ± Week-6 测试集 MAE', 0),
        ('六个标签页：价格趋势 / 敏感性 / 极端场景 / 性能指标 / 误差范围 / 实时数据', 0),
        ('价格趋势：2018–2024 历史重定价 + 实时点；性能指标：Week-6 指标图表', 0),
        ('实时数据：data_updater + GitHub Actions 工作日自动刷新', 0),
    ], top=1.45)
    _add_picture(s, ASSETS / 'ui_trend.png', 6.6, 1.9, width=6.4)
    _footer(s, 9, total); slides.append(s)

    # ── 10. tool screenshots ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '定价工具界面', 'Tool UI (captured from the running app)')
    _add_picture(s, ASSETS / 'ui_overview.png', 0.6, 1.5, width=6.2)
    _add_picture(s, ASSETS / 'ui_metrics.png', 6.9, 1.5, width=6.0)
    _footer(s, 10, total); slides.append(s)

    # ── 11. key findings ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '核心发现与投资启示', 'Key Findings & Investment Insights')
    _bullets(s, [
        ('以市场隐含波动率为目标是定价改善最大来源（VIX-proxy，-23% MAE）', 0),
        ('正则化越过 persistence 门槛（GBDT-anchored 7.16%），实盘三目标全达成', 0),
        ('ML 在极端场景放大波动率风险敞口，但通过特征重工程捕获 VIX 跳升', 0),
        ('投资启示 1：静态 BSM 系统性低估高波动/恐慌期 → 用 ML 定价捕捉 IV premium', 0),
        ('投资启示 2：深度实值 Chooser 波动率不敏感 → 对冲关注 delta 而非 vega', 0),
        ('投资启示 3：双轨定价 + 误差条量化模型不确定性，支持交易决策', 0),
    ], top=1.5)
    _footer(s, 11, total); slides.append(s)

    # ── 12. deliverables ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '交付物清单', 'Deliverables')
    _metric_table(s, [
        ['交付物', '文件', '说明'],
        ['定价工具', 'streamlit_app.py + dashboard.py', '双轨定价 + 误差条 + 完整仪表板'],
        ['最终报告', 'Week_8_实验报告.pdf', '11 页，整合 8 周发现'],
        ['演示文稿', 'Week_8_Final_Presentation.pptx', '本套件'],
        ['演示脚本', 'demo_script.md', '5–10 分钟工具演示视频脚本'],
        ['GitHub 仓库', 'README.md + 代码', '可部署工具，含运行说明'],
    ], top=1.8, col_w=[2.2, 4.4, 5.5])
    _footer(s, 12, total); slides.append(s)

    # ── 13. limitations & future ────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, WHITE)
    _title_box(s, '局限性与未来工作', 'Limitations & Future Work')
    _bullets(s, [
        ('真实 JPM 单票期权历史 IV 不可得（公开 API 仅单日快照）→ VIX 代理延续', 0),
        ('方法二端到端样本瓶颈；接入真实期权价格标签后可再验证', 0),
        ('实盘指标仅单快照（2026-07-27）验证，需真实环境多日回测', 0),
        ('未来：复杂选择权（复合/部分）、随机波动率（Hull-White）、跳跃扩散', 0),
    ], top=1.5)
    _footer(s, 13, total); slides.append(s)

    # ── 14. thank you ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _set_bg(s, DARK)
    box = s.shapes.add_textbox(Inches(1.2), Inches(2.7), Inches(11), Inches(2.0))
    tf = box.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '谢谢 / Thank You'
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = 'Chooser Option Pricing with Real-World Data & ML — 8 周完整交付'
    r2.font.size = Pt(16); r2.font.color.rgb = AMBER
    slides.append(s)

    out = cfg.PRESENTATION_PPTX
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f'saved {out} ({len(slides)} slides)')


if __name__ == '__main__':
    build()
